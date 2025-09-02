from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_minimal_template(path: str):
    """Create a minimal PowerPoint template with one title slide and one content slide.

    Args:
        path (str): The file path to save the template.
    """
    prs = Presentation()

    # Title slide layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Title Placeholder"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Subtitle Placeholder"

    # Content slide layout
    content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Content Title"
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Bullet 1\n Bullet 2"

    # Save template
    prs.save(path)

if __name__ == "__main__":
    out_path = os.path.join(os.path.dirname(__file__), "templates", "base_template.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    create_minimal_template(out_path)
    print(f"Template saved to {out_path}")
