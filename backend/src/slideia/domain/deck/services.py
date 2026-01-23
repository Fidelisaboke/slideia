from pptx import Presentation
from slideia.core.logging import get_logger
from slideia.domain.deck.models import Deck, Slide
from slideia.infra.cache import Cache, RedisCache
from slideia.infra.openrouter import OpenRouterLLM

logger = get_logger(__name__)


def create_minimal_template(path: str):
    """Create a minimal PowerPoint template with one title slide and one content slide.

    Args:
        path (str): The file path to save the template.
    """
    prs = Presentation()

    # Title slide layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Title Placeholder"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Subtitle Placeholder"

    # Content slide layout
    content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Content Title"
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Bullet 1\n Bullet 2"

    # Save template
    prs.save(path)


def generate_full_deck(
    topic: str,
    audience: str,
    tone: str,
    slide_count: int,
    llm: OpenRouterLLM,
    cache: Cache | RedisCache,
) -> Deck:
    cached = cache.get(topic, audience, tone, slide_count)
    if cached:
        return Deck(
            outline=cached.get("outline", {}),
            slides=[Slide(**s) for s in cached.get("slides", [])],
        )

    logger.info("Generating new deck...")

    outline = llm.propose_outline(
        topic=topic,
        audience=audience,
        tone=tone,
        slide_count=slide_count,
    )

    slides_content = []
    for slide_spec in outline.get("slides", []):
        slides_content.append(llm.draft_slide(slide_spec))

    logger.info("Deck generation complete!")
    result = {
        "outline": outline,
        "slides": slides_content,
    }

    cache.set(topic, audience, tone, slide_count, result)
    logger.info("Cached the generated deck.")

    return Deck(
        outline=outline,
        slides=[Slide(**s) for s in slides_content],
    )
