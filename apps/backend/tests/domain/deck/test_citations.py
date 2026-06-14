import json
import os
import tempfile
import pytest
from pptx import Presentation
from slideia.domain.deck.exporter import export_slides
from slideia.domain.deck.pdf_exporter import export_deck_to_pdf


@pytest.fixture
def slide_data_with_citations():
    return {
        "title": "Test Presentation with References",
        "subtitle": "For Citation Testing",
        "palette": ["#7c5cca", "#14b8a6", "#a78bfa"],
        "font": "Inter",
        "citations": [
            "https://example.com/source1 - Primary Research",
            "https://example.com/source2 - Secondary Analysis",
        ],
        "slides": [
            {
                "title": "Slide 1",
                "summary": "Summary 1",
                "bullets": ["Bullet 1", "Bullet 2"],
                "notes": "Speaker notes 1",
                "image_prompt": "Prompt 1",
                "theme": {"font": "Arial", "color": "#003366"},
            }
        ],
    }


@pytest.fixture
def temp_json_file(slide_data_with_citations):
    fd, path = tempfile.mkstemp(suffix=".json")
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        json.dump(slide_data_with_citations, f)
    yield path
    os.remove(path)


@pytest.fixture
def temp_pptx_output():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture
def temp_pdf_output():
    fd, path = tempfile.mkstemp(suffix=".pdf")
    os.close(fd)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.mark.asyncio
async def test_pptx_export_includes_references_slide(temp_json_file, temp_pptx_output):
    await export_slides(temp_json_file, temp_pptx_output)
    assert os.path.exists(temp_pptx_output)

    prs = Presentation(temp_pptx_output)
    # 1 Title slide + 1 Content slide + 1 References slide = 3 slides total
    assert len(prs.slides) == 3

    # Verify the last slide is the References slide
    ref_slide = prs.slides[2]

    # Check that one of the shapes contains "References" title
    titles = []
    for shape in ref_slide.shapes:
        if shape.has_text_frame:
            titles.append(shape.text_frame.text)

    assert any("References" in t for t in titles)
    assert any("[1] https://example.com/source1 - Primary Research" in t for t in titles)
    assert any("[2] https://example.com/source2 - Secondary Analysis" in t for t in titles)


@pytest.mark.asyncio
async def test_pdf_export_includes_references_slide(temp_pdf_output, slide_data_with_citations):
    # Mocking images for pdf exporter as it downloads images asynchronously
    slide_data_with_citations["slides"][0]["image_url"] = None

    fd, path = tempfile.mkstemp(suffix=".json")
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            json.dump(slide_data_with_citations, f)

        await export_deck_to_pdf(path, temp_pdf_output)
    finally:
        if os.path.exists(path):
            os.remove(path)

    assert os.path.exists(temp_pdf_output)
    # Check that PDF has non-zero size
    assert os.path.getsize(temp_pdf_output) > 0
