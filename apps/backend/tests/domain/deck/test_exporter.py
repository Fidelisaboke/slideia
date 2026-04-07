import json
import os
import tempfile

import pytest
from pptx import Presentation
from slideia.domain.deck.exporter import export_slides


@pytest.fixture
def minimal_json_slide():
    return {
        "title": "Test Presentation",
        "subtitle": "For Testing",
        "slides": [
            {
                "title": "Slide 1",
                "summary": "Summary 1",
                "bullets": ["Bullet 1", "Bullet 2"],
                "notes": "Speaker notes 1",
                "image_prompt": "Prompt 1",
                "theme": {"font": "Arial", "color": "#003366"},
            },
            {
                "title": "Slide 2",
                "summary": "Summary 2",
                "bullets": ["Bullet 3"],
                "notes": "Speaker notes 2",
                "image_prompt": "Prompt 2",
                "theme": {"font": "Calibri", "color": "#222222"},
            },
        ],
    }


@pytest.fixture
def temp_json_file(minimal_json_slide):
    fd, path = tempfile.mkstemp(suffix=".json")
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        json.dump(minimal_json_slide, f)
    yield path
    os.remove(path)


@pytest.fixture
def temp_output_file():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture(scope="module")
def clean_templates():
    # Remove generated template if it exists
    template_dir = os.path.join(
        os.path.dirname(__file__), "../../../../src/slideia/domain/deck/templates"
    )
    template_path = os.path.join(template_dir, "base_template.pptx")
    if os.path.exists(template_path):
        os.remove(template_path)
    yield
    if os.path.exists(template_path):
        os.remove(template_path)


def test_export_slides_creates_pptx(temp_json_file, temp_output_file, clean_templates):
    export_slides(temp_json_file, temp_output_file)
    assert os.path.exists(temp_output_file)
    prs = Presentation(temp_output_file)
    # Should have 3 slides: 1 title + 2 content
    assert len(prs.slides) == 3
    # Check title slide
    title_slide = prs.slides[0]
    assert title_slide.shapes.title.text == "Test Presentation"
    # Check content slide titles
    assert prs.slides[1].shapes[0].text_frame.text == "Slide 1"
    assert prs.slides[2].shapes[0].text_frame.text == "Slide 2"


def test_export_slides_missing_input_raises(temp_output_file):
    with pytest.raises(FileNotFoundError):
        export_slides("/nonexistent/file.json", temp_output_file)


def test_export_slides_missing_template_creates_it(
    temp_json_file, temp_output_file, clean_templates
):
    # Remove template if exists
    # Match the template path logic in exporter.py
    exporter_dir = os.path.abspath(
        os.path.join(os.path.dirname(__file__), "../../../src/slideia/domain/deck")
    )
    template_dir = os.path.join(exporter_dir, "templates")
    template_path = os.path.join(template_dir, "base_template.pptx")
    if os.path.exists(template_path):
        os.remove(template_path)
    export_slides(temp_json_file, temp_output_file)
    assert os.path.exists(template_path)
    assert os.path.exists(temp_output_file)


def test_export_slides_handles_type_mismatches(temp_output_file, clean_templates):
    # theme as string, bullets as string, title as int, notes as int
    bad_slide = {
        "title": 123,
        "subtitle": None,
        "slides": [
            {
                "title": 456,
                "summary": 789,
                "bullets": "Bullet as string\nAnother bullet",
                "notes": 1011,
                "image_prompt": 1213,
                "theme": "notadict",
            }
        ],
    }
    fd, path = tempfile.mkstemp(suffix=".json")
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        json.dump(bad_slide, f)
    export_slides(path, temp_output_file)
    prs = Presentation(temp_output_file)
    # Should have 2 slides: 1 title + 1 content
    assert len(prs.slides) == 2
    # Title slide title should be str(123)
    assert prs.slides[0].shapes.title.text == "123"
    # Content slide title should be str(456)
    assert prs.slides[1].shapes[0].text_frame.text == "456"
    os.remove(path)


def test_export_slides_empty_slides(temp_output_file, clean_templates):
    empty_slide = {"title": "Empty", "slides": []}
    fd, path = tempfile.mkstemp(suffix=".json")
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        json.dump(empty_slide, f)
    export_slides(path, temp_output_file)
    prs = Presentation(temp_output_file)
    # Should have 1 slide (title only)
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == "Empty"
    os.remove(path)
