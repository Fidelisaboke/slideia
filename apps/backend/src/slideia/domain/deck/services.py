import asyncio
from typing import AsyncGenerator
from pptx import Presentation
from slideia.core.config import settings
from slideia.core.logging import get_logger
from slideia.domain.deck.models import Deck, Slide
from slideia.infra.cache import Cache, RedisCache
from slideia.infra.openrouter import OpenRouterLLM

logger = get_logger(__name__)


def create_minimal_template(path: str):
    """Create a minimal PowerPoint template with one title slide and one content slide.

    Args:
        path (str): The file path to save the template.
    """
    prs = Presentation()

    # Title slide layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Title Placeholder"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Subtitle Placeholder"

    # Content slide layout
    content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Content Title"
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Bullet 1\n Bullet 2"

    # Save template
    prs.save(path)


async def generate_full_deck(
    topic: str,
    audience: str,
    tone: str,
    slide_count: int,
    llm: OpenRouterLLM,
    cache: Cache | RedisCache,
) -> Deck:
    cached = cache.get(topic, audience, tone, slide_count)
    if cached:
        return Deck(
            outline=cached.get("outline", {}),
            slides=[Slide(**s) for s in cached.get("slides", [])],
        )

    logger.info("Generating new deck...")

    outline = await llm.propose_outline(
        topic=topic,
        audience=audience,
        tone=tone,
        slide_count=slide_count,
    )

    # Use Semaphore to limit concurrent calls to respect rate limits
    semaphore = asyncio.Semaphore(settings.MAX_CONCURRENT_LLM_CALLS)
    batch_size = 3
    slide_specs = outline.get("slides", [])
    batches = [slide_specs[i : i + batch_size] for i in range(0, len(slide_specs), batch_size)]

    async def process_batch(batch):
        async with semaphore:
            result = await llm.draft_slides_batch(topic, audience, batch)
            return result.get("slides", [])

    logger.info(f"Drafting {len(slide_specs)} slides in {len(batches)} batches...")
    tasks = [process_batch(b) for b in batches]
    results = await asyncio.gather(*tasks)

    # Flatten results
    slides_content = [slide for batch_result in results for slide in batch_result]

    logger.info("Deck generation complete!")
    result = {
        "outline": outline,
        "slides": slides_content,
    }

    cache.set(topic, audience, tone, slide_count, result)
    logger.info("Cached the generated deck.")

    return Deck(
        outline=outline,
        slides=[Slide(**s) for s in slides_content],
    )


async def generate_full_deck_stream(
    topic: str,
    audience: str,
    tone: str,
    slide_count: int,
    llm: OpenRouterLLM,
    cache: Cache | RedisCache,
) -> AsyncGenerator[dict, None]:
    """
    Generate a full deck and yield progress events.
    """
    cached = cache.get(topic, audience, tone, slide_count)
    if cached:
        logger.info("Serving cached deck via stream.")
        yield {"step": "complete", "progress": 100, "message": "Loaded from cache!", "data": cached}
        return

    logger.info("Starting streaming generation...")

    # Step 1: Outline
    yield {"step": "outline", "progress": 10, "message": "Analyzing topic and structuring the story..."}

    outline = await llm.propose_outline(
        topic=topic,
        audience=audience,
        tone=tone,
        slide_count=slide_count,
    )

    slide_specs = outline.get("slides", [])
    total_slides = len(slide_specs)
    slides_content = [None] * total_slides
    slides_processed = 0

    # Step 2: Slides (Batched and Parallel)
    semaphore = asyncio.Semaphore(settings.MAX_CONCURRENT_LLM_CALLS)
    batch_size = 3
    batches = [slide_specs[i : i + batch_size] for i in range(0, len(slide_specs), batch_size)]

    async def process_batch_with_progress(batch, start_idx):
        async with semaphore:
            result = await llm.draft_slides_batch(topic, audience, batch)
            return result.get("slides", []), start_idx

    tasks = [process_batch_with_progress(b, i * batch_size) for i, b in enumerate(batches)]

    for future in asyncio.as_completed(tasks):
        batch_slides, batch_start = await future
        for j, slide in enumerate(batch_slides):
            idx = batch_start + j
            if idx < total_slides:
                slides_content[idx] = slide
                slides_processed += 1
                progress = 10 + int((slides_processed / total_slides) * 80)

                yield {
                    "step": "slide",
                    "index": idx + 1,
                    "total": total_slides,
                    "title": slide.get("title", "Untitled"),
                    "progress": progress,
                    "message": f"Drafted slide {slides_processed} of {total_slides}: {slide.get('title')}",
                }

    # Ensure all slots are filled (in case of model errors we could have None)
    slides_content = [s for s in slides_content if s is not None]

    # Step 3: Complete
    result = {
        "outline": outline,
        "slides": slides_content,
    }

    cache.set(topic, audience, tone, slide_count, result)

    yield {"step": "complete", "progress": 100, "message": "Presentation ready!", "data": result}


async def propose_outline_stream(
    topic: str,
    audience: str,
    tone: str,
    slide_count: int,
    llm: OpenRouterLLM,
    cache: Cache | RedisCache,
) -> AsyncGenerator[dict, None]:
    """
    Propose an outline and yield progress events.
    """
    # Check cache
    cached = cache.get(topic, audience, tone, slide_count)
    if cached:
        logger.info("Serving cached outline via stream.")
        yield {
            "step": "complete",
            "progress": 100,
            "message": "Loaded from cache!",
            "data": cached.get("outline", {}),
        }
        return

    logger.info("Starting streaming outline proposal...")

    yield {"step": "outline", "progress": 20, "message": "Analyzing theme and context..."}

    # We don't have sub-steps for LLM.propose_outline yet,
    # but we can simulate progress and use the async call.
    yield {"step": "outline", "progress": 50, "message": "Drafting structural framework..."}

    outline = await llm.propose_outline(
        topic=topic,
        audience=audience,
        tone=tone,
        slide_count=slide_count,
    )

    yield {"step": "outline", "progress": 90, "message": "Finalizing presentation structure..."}

    # We yield the final result
    yield {"step": "complete", "progress": 100, "message": "Outline complete!", "data": outline}
