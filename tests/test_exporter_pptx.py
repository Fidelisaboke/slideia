"""
Test for exporter: verifies .pptx generation and slide count.
"""
import os
import pytest

from pptx import Presentation

from slideia.tools.exporter import export_slides

@pytest.fixture(scope="session", autouse=True)
def ensure_template_exists():
    template_dir = os.path.join(os.path.dirname(__file__), "..", "src", "slideia", "tools", "templates")
    os.makedirs(template_dir, exist_ok=True)
    template_path = os.path.join(template_dir, "base_template.pptx")
    if not os.path.exists(template_path):
        prs = Presentation()
        prs.save(template_path)


def test_exporter_creates_pptx(sample_deck_json, tmp_export_dir):
    output_path = os.path.join(tmp_export_dir, "test_deck.pptx")
    export_slides(sample_deck_json, output_path)
    assert os.path.exists(output_path)
    prs = Presentation(output_path)
    # 1 title + 2 content slides
    assert len(prs.slides) == 3


def test_exporter_handles_missing_image(tmp_export_dir, sample_deck_json):
    # Remove image_path from sample, but keep image_prompt
    import json
    with open(sample_deck_json, "r", encoding="utf-8") as f:
        deck = json.load(f)
    deck["slides"][0]["image_path"] = "/nonexistent/path.png"
    new_json = os.path.join(tmp_export_dir, "deck_missing_img.json")
    with open(new_json, "w", encoding="utf-8") as f:
        json.dump(deck, f)
    output_path = os.path.join(tmp_export_dir, "test_missing_img.pptx")
    export_slides(new_json, output_path)
    assert os.path.exists(output_path)
    prs = Presentation(output_path)
    assert len(prs.slides) == 3
